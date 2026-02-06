import os
import json
from pathlib import Path

# 1. 强制重定向缓存目录 (必须在 import transformers 之前!)
os.environ["HF_HOME"] = "./hf_cache"

import jieba
import torch
import numpy as np
from typing import List, Dict, Union
from rouge_chinese import Rouge
from transformers import AutoModelForCausalLM, AutoTokenizer

# 引入处理文件的库
from pptx import Presentation
from pypdf import PdfReader

class MetricCalculator:
    def __init__(self, ppl_model_path: str = None, hf_token: str = None, device: str = "cuda"):
        self.rouge = Rouge()
        self.device = device if torch.cuda.is_available() else "cpu"
        self.ppl_model = None
        self.ppl_tokenizer = None
        self.ppl_model_path = ppl_model_path
        self.hf_token = hf_token

    def _load_ppl_model(self):
        if self.ppl_model is None:
            print(f"Loading PPL model from {self.ppl_model_path}...")
            try:
                self.ppl_tokenizer = AutoTokenizer.from_pretrained(
                    self.ppl_model_path, 
                    token=self.hf_token
                )
                
                if self.ppl_tokenizer.pad_token is None:
                    self.ppl_tokenizer.pad_token = self.ppl_tokenizer.eos_token
                
                self.ppl_model = AutoModelForCausalLM.from_pretrained(
                    self.ppl_model_path, 
                    device_map="auto",  
                    torch_dtype=torch.float16, 
                    token=self.hf_token
                )
                self.ppl_model.eval()
            except Exception as e:
                raise RuntimeError(f"模型加载失败: {e}")

    def calculate_rouge_l(self, hypothesis: str, reference: str) -> float:
        """
        计算 ROUGE-L
        hypothesis: PPT内容 (生成的摘要)
        reference: PDF内容 (原文)
        """
        # 使用 jieba 分词
        hyp_tokens = " ".join(jieba.cut(str(hypothesis)))
        ref_tokens = " ".join(jieba.cut(str(reference)))
        
        if not hyp_tokens.strip() or not ref_tokens.strip(): 
            return 0.0
            
        try:
            scores = self.rouge.get_scores(hyp_tokens, ref_tokens)
            return scores[0]["rouge-l"]["f"]
        except Exception as e:
            print(f"ROUGE计算出错: {e}")
            return 0.0

    def calculate_ppl(self, texts: List[str]) -> float:
        """
        计算 PPL
        texts: PPT中的每一页或每一段文本组成的列表
        """
        self._load_ppl_model()
        ppls = []
        
        # 过滤空字符串
        valid_texts = [t for t in texts if t.strip()]
        if not valid_texts:
            return 0.0

        print(f"正在计算 PPL (共 {len(valid_texts)} 段文本)...")
        
        for text in valid_texts:
            try:
                inputs = self.ppl_tokenizer(text, return_tensors="pt").to(self.ppl_model.device)
                
                # 截断过长的输入
                if inputs.input_ids.size(1) > self.ppl_model.config.max_position_embeddings:
                    inputs.input_ids = inputs.input_ids[:, :self.ppl_model.config.max_position_embeddings]
                    
                with torch.no_grad():
                    outputs = self.ppl_model(inputs.input_ids, labels=inputs.input_ids)
                    loss = outputs.loss
                    perplexity = torch.exp(loss)
                    
                    if not torch.isnan(perplexity): 
                        ppls.append(perplexity.item())
            except Exception as e:
                print(f"计算某段文本PPL时出错: {e}")
                continue

        if not ppls: 
            return 0.0
        return sum(ppls) / len(ppls)

# ================= 工具函数：提取文本 =================

def extract_text_from_pptx(pptx_path: str) -> (List[str], str):
    """
    提取PPT内容。
    返回: 
    1. List[str]: 用于计算 PPL (保留每页/每块的结构)
    2. str: 用于计算 ROUGE (全文合并)
    """
    if not os.path.exists(pptx_path):
        return [], ""
    
    prs = Presentation(pptx_path)
    slides_text_list = []
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        # 将每一页的内容拼成一个字符串，作为一个 PPL 计算单元
        if slide_text:
            slides_text_list.append("\n".join(slide_text))
            
    full_text = "\n".join(slides_text_list)
    return slides_text_list, full_text

def extract_text_from_pdf(pdf_path: str) -> str:
    """提取PDF全文内容用于作为 ROUGE 的 Reference"""
    if not os.path.exists(pdf_path):
        return ""
    
    try:
        reader = PdfReader(pdf_path)
        text_content = []
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text_content.append(extracted)
        return "\n".join(text_content)
    except Exception as e:
        print(f"读取PDF失败 {pdf_path}: {e}")
        return ""

# ================= 核心逻辑：评测单个文件夹 =================

def evaluate_paper_folder(folder_path: str, calculator: MetricCalculator):
    """
    评测指定文件夹下的 PPTX 和 PDF (自动查找 PDF)
    """
    folder_path = Path(folder_path)
    pptx_path = folder_path / "merged_presentation.pptx"
    result_path = folder_path / "eval_result_exist_metric.json"

    print(f"\n>>> 正在评测文件夹: {folder_path}")

    # ================= 修改开始 =================
    # 1. 自动查找当前目录下的 PDF 文件
    pdf_files = list(folder_path.glob("*.pdf"))
    
    if not pdf_files:
        print(f"Error: 在 {folder_path} 下未找到任何 PDF 文件，跳过。")
        return
    
    # 默认取第一个找到的 PDF，如果有多个可以加个警告
    pdf_path = pdf_files[0]
    if len(pdf_files) > 1:
        print(f"Warning: 该目录下发现多个 PDF，默认使用第一个: {pdf_path.name}")
    
    print(f"已锁定 PDF 文件: {pdf_path.name}")
    # ================= 修改结束 =================

    # 2. 检查 PPTX 是否存在
    if not pptx_path.exists():
        print(f"Error: 找不到 {pptx_path}，跳过。")
        return

    # 3. 提取内容
    print("提取 PPTX 内容...")
    pptx_slides_list, pptx_full_text = extract_text_from_pptx(str(pptx_path))
    
    print("提取 PDF 内容...")
    pdf_full_text = extract_text_from_pdf(str(pdf_path))

    if not pptx_full_text or not pdf_full_text:
        print("Error: 提取的文本内容为空，跳过计算。")
        return

    # 4. 计算指标
    results = {}
    
    # ROUGE-L (PPT全文 vs PDF全文)
    print("计算 ROUGE-L...")
    rouge_score = calculator.calculate_rouge_l(pptx_full_text, pdf_full_text)
    results["rouge_l"] = rouge_score * 100  # 转为百分制
    results["pdf_filename"] = pdf_path.name  # 顺便记录一下用的是哪个PDF
    print(f"ROUGE-L: {rouge_score * 100:.4f}")

    # PPL (PPT List)
    print("计算 PPL...")
    try:
        ppl_score = calculator.calculate_ppl(pptx_slides_list)
        results["ppl"] = ppl_score
        print(f"PPL: {ppl_score:.4f}")
    except Exception as e:
        print(f"PPL计算失败: {e}")
        results["ppl"] = -1

    # 5. 保存结果
    try:
        with open(result_path, "w", encoding="utf-8") as f:
            json.dump(results, f, ensure_ascii=False, indent=4)
        print(f"结果已保存至: {result_path}")
    except Exception as e:
        print(f"保存结果失败: {e}")


import argparse


def main():
    # 1. 配置命令行参数解析
    parser = argparse.ArgumentParser(description="LLM 评测工具 - 支持单文件夹或批量模式")
    
    # 位置参数：目标路径
    parser.add_argument("input_path", type=str, help="指定论文文件夹路径或批量文件根目录路径")
    
    # 可选参数：是否开启批量模式
    parser.add_argument("--batch", action="store_true", help="启用批量模式，处理输入路径下的所有子文件夹")
    
    # 模型配置参数 (也可以通过命令行传入，这里先保留你的默认值)
    parser.add_argument("--model", type=str, default="meta-llama/Meta-Llama-3-8B", help="HF 模型路径")
    parser.add_argument("--token", type=str, help="HuggingFace Token")

    args = parser.parse_args()

    # 2. 初始化计算器 (只加载一次模型)
    print(f"正在加载模型: {args.model} ...")
    calculator = MetricCalculator(
        ppl_model_path=args.model, 
        hf_token=args.token
    )

    target_path = Path(args.input_path)
    
    if not target_path.exists():
        print(f"错误: 路径 {target_path} 不存在")
        return

    # 3. 根据模式执行逻辑
    if args.batch:
        # 批量模式：处理子目录
        subdirs = [d for d in target_path.iterdir() if d.is_dir()]
        print(f"【批量模式】在 {target_path} 下发现 {len(subdirs)} 个文件夹。")
        
        for idx, subdir in enumerate(subdirs):
            print(f"\n进度: {idx+1}/{len(subdirs)} | 正在处理: {subdir.name}")
            evaluate_paper_folder(str(subdir), calculator)
    else:
        # 单文件夹模式
        print(f"【单文件夹模式】正在处理: {target_path.name}")
        evaluate_paper_folder(str(target_path), calculator)

    print("\n所有任务已完成。")

if __name__ == "__main__":
    main()
