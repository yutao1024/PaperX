import os
import time
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn  # 必须引入

from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By


def init_driver():
    """初始化无头 Chrome 浏览器"""
    chrome_options = Options()
    chrome_options.add_argument("--headless")
    chrome_options.add_argument("--disable-gpu")
    chrome_options.add_argument("--no-sandbox")
    chrome_options.add_argument("--hide-scrollbars")
    chrome_options.add_argument("--window-size=1920,1080")
    chrome_options.add_argument("--allow-file-access-from-files")
    
    service = Service(ChromeDriverManager().install())
    driver = webdriver.Chrome(service=service, options=chrome_options)
    return driver

def get_html_files(input_dir):
    """获取目录下所有 html 文件并排序"""
    files = [f for f in os.listdir(input_dir) if f.lower().endswith('.html')]
    
    def sort_key(filename):
        try:
            if '_ppt' in filename:
                number_part = filename.split('_ppt')[0]
                return int(number_part)
            else:
                digits = ''.join(filter(str.isdigit, filename))
                return int(digits) if digits else float('inf')
        except:
            return float('inf')

    files.sort(key=sort_key)
    return files

def set_font_transparency(font, transparency_percent=100):
    """
    设置字体透明度的辅助函数 (修复版)
    """
    # OpenXML 中 alpha 值范围是 0 (全透明) 到 100000 (不透明)
    alpha_val = int((100 - transparency_percent) * 1000)
    
    # 1. 必须先设置为纯色填充
    font.fill.solid()
    # 2. 设置基底颜色 (黑色)，这会创建 <a:srgbClr> 节点
    font.fill.fore_color.rgb = RGBColor(0, 0, 0) 
    
    # 3. 获取底层 XML 元素 (_xFill 就是 <a:solidFill>)
    fill = font.fill.fore_color._xFill
    
    # 4. 直接访问 srgbClr (之前错误的写法是 fill.solidFill.srgbClr)
    if hasattr(fill, 'srgbClr') and fill.srgbClr is not None:
        # 移除已有的 alpha 节点
        for alpha in fill.srgbClr.findall(qn('a:alpha')):
            fill.srgbClr.remove(alpha)
        
        # 创建并添加新的 alpha 节点
        alpha_element = fill.srgbClr.makeelement(qn('a:alpha'))
        alpha_element.set('val', str(alpha_val))
        fill.srgbClr.append(alpha_element)

def convert_to_pptx(input_dir, output_dir, output_filename="merged_presentation.pptx"):
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    
    if not output_path.exists():
        os.makedirs(output_path)

    html_files = get_html_files(input_path)
    if not html_files:
        print(f"在目录 {input_dir} 中未找到 HTML 文件。")
        return

    print(f"找到 {len(html_files)} 个 HTML 文件")

    # 1. 初始化 PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # 2. 启动浏览器
    driver = init_driver()

    try:
        for filename in html_files:
            file_path = input_path / filename
            abs_file_path = file_path.resolve().as_uri()
            
            print(f"正在处理: {filename} ...")
            driver.get(abs_file_path)
            
            # CSS 注入
            driver.execute_script("""
                document.body.style.margin = '0';
                document.body.style.padding = '0';
                document.body.style.overflow = 'hidden';
                var slides = document.getElementsByClassName('slide');
                for (var i = 0; i < slides.length; i++) {
                    slides[i].style.display = 'inline-block';
                    slides[i].style.width = 'fit-content';
                }
            """)
            
            time.sleep(0.5)

            # --- 提取文本 ---
            try:
                text_content = driver.find_element(By.TAG_NAME, "body").text
            except Exception as e:
                print(f"  警告: 文本提取失败 - {e}")
                text_content = ""

            # --- 截图逻辑 ---
            try:
                element = driver.find_element(By.TAG_NAME, "svg")
            except:
                try:
                    element = driver.find_element(By.CLASS_NAME, "slide")
                except:
                    element = driver.find_element(By.TAG_NAME, "body")

            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                element.screenshot(tmp_img.name)
                tmp_img_path = tmp_img.name

            # 3. 添加幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # 空白页

            # 插入图片
            slide.shapes.add_picture(
                tmp_img_path, 
                0, 
                0, 
                width=prs.slide_width, 
                height=prs.slide_height
            )

            # --- 插入透明文本框 ---
            if text_content:
                txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
                tf = txBox.text_frame
                tf.word_wrap = True 
                
                p = tf.paragraphs[0]
                p.text = text_content
                p.font.size = Pt(10)
                
                # 设置字体全透明
                try:
                    set_font_transparency(p.font, transparency_percent=100)
                except Exception as e:
                    print(f"  设置透明度失败 (已忽略): {e}")

            try:
                os.remove(tmp_img_path)
            except:
                pass

        # 4. 保存
        output_file_path = output_path / output_filename
        prs.save(output_file_path)
        print(f"\n成功! PPTX 文件已生成: {output_file_path}")

    except Exception as e:
        print(f"发生错误: {e}")
        import traceback
        traceback.print_exc()
    finally:
        driver.quit()

# if __name__ == "__main__":
#     # 请根据实际情况修改目录
#     INPUT_DIRECTORY = r"/home/yutao/agent/PPTAgent/html2pptx/html_input/HSCR_Hierarchical_Self-Contrastive_Rewarding_for_Aligning_Medical_Vision_Language_Models"
#     OUTPUT_DIRECTORY = r"/home/yutao/agent/PPTAgent/html2pptx/pptx_output/HSCR_Hierarchical_Self-Contrastive_Rewarding_for_Aligning_Medical_Vision_Language_Models_havetext"
    
#     convert_to_pptx(INPUT_DIRECTORY, OUTPUT_DIRECTORY)